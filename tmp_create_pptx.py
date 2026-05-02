from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# カラー定義
COL_DARK = RGBColor(0x1A, 0x1A, 0x2E)
COL_BLUE = RGBColor(0x16, 0x78, 0xC2)
COL_GREEN = RGBColor(0x27, 0xAE, 0x60)
COL_PURPLE = RGBColor(0x8E, 0x44, 0xAD)
COL_ACCENT = RGBColor(0xF3, 0x9C, 0x12)
COL_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COL_LIGHT = RGBColor(0xF5, 0xF7, 0xFA)
COL_TEXT = RGBColor(0x2C, 0x3E, 0x50)
COL_GRAY = RGBColor(0x95, 0xA5, 0xA6)
COL_C = RGBColor(0x16, 0x78, 0xC2)
COL_A = RGBColor(0x27, 0xAE, 0x60)
COL_D = RGBColor(0xE6, 0x7E, 0x22)
COL_E = RGBColor(0x8E, 0x44, 0xAD)

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)


def add_rect(slide, l, t, w, h, fill_color=None, line_color=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, l, t, w, h, font_size=14, bold=False, color=None, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return txBox


def add_multiline(slide, lines, l, t, w, h, font_size=10, bold=False, color=None, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.bold = bold
        if color:
            run.font.color.rgb = color
    return txBox


def slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


blank_layout = prs.slide_layouts[6]
col_positions = [0.25, 3.55, 6.85, 10.05]

# ========== スライド1: 表紙 ==========
slide = prs.slides.add_slide(blank_layout)
slide_bg(slide, COL_DARK)
add_rect(slide, 0, 0, 13.33, 0.08, COL_BLUE)
add_text(slide, 'CX推進部', 1.5, 1.8, 10, 0.8, font_size=20, color=COL_GRAY, align=PP_ALIGN.CENTER)
add_text(slide, 'DX推進ロードマップ', 1.0, 2.5, 11.33, 1.2, font_size=44, bold=True, color=COL_WHITE, align=PP_ALIGN.CENTER)
add_text(slide, '2026年4月 | 経営層報告・チーム実行計画 兼用版', 1.5, 3.9, 10, 0.6, font_size=16, color=COL_GRAY, align=PP_ALIGN.CENTER)
phases = [
    ('DX 1.0', '〜2028年3月', COL_BLUE, 1.0),
    ('DX 1.5', '2028〜2031年', COL_GREEN, 5.0),
    ('DX 2.0', '2031年〜', COL_PURPLE, 9.0),
]
for name, period, col, x in phases:
    add_rect(slide, x, 5.0, 3.4, 0.7, col)
    add_text(slide, name, x+0.1, 5.05, 1.5, 0.4, font_size=14, bold=True, color=COL_WHITE)
    add_text(slide, period, x+0.1, 5.42, 3.2, 0.3, font_size=9, color=COL_WHITE)
add_rect(slide, 0, 7.35, 13.33, 0.08, COL_ACCENT)
add_text(slide, '杏林堂薬局 CX推進部', 0.3, 7.1, 5, 0.3, font_size=9, color=COL_GRAY)
print('スライド1完了')

# ========== スライド2: ビジョンとフェーズ概観 ==========
slide = prs.slides.add_slide(blank_layout)
slide_bg(slide, COL_LIGHT)
add_rect(slide, 0, 0, 13.33, 0.6, COL_DARK)
add_rect(slide, 0, 0.6, 13.33, 0.05, COL_BLUE)
add_text(slide, 'DXビジョンとフェーズ全体像', 0.3, 0.1, 10, 0.45, font_size=20, bold=True, color=COL_WHITE)

add_rect(slide, 0.3, 0.85, 12.7, 0.85, RGBColor(0xEB, 0xF5, 0xFB))
add_rect(slide, 0.3, 0.85, 0.08, 0.85, COL_BLUE)
add_text(slide, 'DXビジョン：デジタルを活用した「伴走型支援」によるCX向上と業務変革', 0.55, 0.92, 11, 0.35, font_size=13, bold=True, color=COL_TEXT)
add_text(slide, '企業理念「健康で豊かな生活への貢献」をDXで実現。CX推進部が店舗・顧客・データをつなぐ価値統合を推進。', 0.55, 1.25, 11.5, 0.35, font_size=10, color=COL_TEXT)

add_rect(slide, 0.5, 2.35, 12.3, 0.04, COL_GRAY)
years = ['2026', '2027', '2028', '2029', '2030', '2031', '2032〜']
for i, yr in enumerate(years):
    x = 0.5 + i * (12.3/6)
    add_rect(slide, x-0.01, 2.25, 0.04, 0.22, COL_GRAY)
    add_text(slide, yr, x-0.3, 2.48, 0.7, 0.3, font_size=8, color=COL_GRAY, align=PP_ALIGN.CENTER)

phase_data = [
    ('DX 1.0', '業務自動化・デジタル基盤整備', COL_BLUE, 0.5, 4.0),
    ('DX 1.5', 'データ活用・分析による意思決定支援', COL_GREEN, 4.6, 4.0),
    ('DX 2.0', 'AI活用・自律最適化・顧客体験革新', COL_PURPLE, 8.55, 4.1),
]
for name, desc, col, x, w in phase_data:
    add_rect(slide, x, 1.85, w, 0.5, col)
    add_text(slide, name, x+0.1, 1.88, w-0.2, 0.25, font_size=13, bold=True, color=COL_WHITE)
    add_text(slide, desc, x+0.1, 2.1, w-0.2, 0.22, font_size=8, color=COL_WHITE)

milestones = [
    (1.52, ['2027/3', 'MS365', '店舗展開'], COL_BLUE),
    (4.6, ['2028', 'DX1.0', '完了'], COL_BLUE),
    (6.63, ['2028', 'ツルハHD', 'DB刷新連携'], COL_GREEN),
    (8.55, ['2030', 'DX1.5', '定着'], COL_GREEN),
    (10.6, ['2031', 'AI人財', '内製化'], COL_PURPLE),
]
for x, label_lines, col in milestones:
    add_rect(slide, x-0.07, 2.2, 0.16, 0.16, col)
    add_multiline(slide, label_lines, x-0.5, 2.78, 1.1, 0.8, font_size=7.5, color=col, align=PP_ALIGN.CENTER)

add_text(slide, '重点テーマの優先順位', 0.3, 3.8, 5, 0.3, font_size=11, bold=True, color=COL_TEXT)
themes = [('C データ基盤', COL_C, 0.3), ('A シフト・人員管理', COL_A, 3.8), ('D 顧客体験', COL_D, 7.4)]
for t, col, x in themes:
    add_rect(slide, x, 4.15, 3.3, 0.5, col)
    add_text(slide, t, x+0.1, 4.22, 3.0, 0.35, font_size=12, bold=True, color=COL_WHITE)

add_rect(slide, 0.3, 4.85, 12.7, 1.85, COL_WHITE)
add_text(slide, 'フェーズ別 主要KPI目標', 0.5, 4.9, 5, 0.3, font_size=10, bold=True, color=COL_TEXT)
kpi_data = [
    ('DX 1.0', '業務時間 20%削減\nシフトシステム全店100%', COL_BLUE, 0.5),
    ('DX 1.5', 'データドリブン意思決定 70%\n定例報告自動化 60%', COL_GREEN, 4.5),
    ('DX 2.0', 'AI支援業務 50%以上\nシフト自動生成 90%', COL_PURPLE, 8.5),
]
for phase, kpi, col, x in kpi_data:
    add_rect(slide, x, 5.25, 3.5, 0.22, col)
    add_text(slide, phase, x+0.05, 5.27, 3.0, 0.18, font_size=9, bold=True, color=COL_WHITE)
    add_multiline(slide, kpi.split('\n'), x, 5.52, 3.5, 0.7, font_size=9, color=COL_TEXT)
print('スライド2完了')

# ========== スライド3: マイルストーン ==========
slide = prs.slides.add_slide(blank_layout)
slide_bg(slide, COL_LIGHT)
add_rect(slide, 0, 0, 13.33, 0.6, COL_DARK)
add_rect(slide, 0, 0.6, 13.33, 0.05, COL_ACCENT)
add_text(slide, '主要マイルストーン', 0.3, 0.1, 10, 0.45, font_size=20, bold=True, color=COL_WHITE)

ms_items = [
    ('2026年度', 'MS365 本部員全員活用開始', 'Teamsによる情報共有・SharePointの整備を推進', COL_BLUE),
    ('2027年3月', 'MS365 店舗PC展開完了', '店舗スタッフへのMS365展開。デジタル化の全社基盤が整う', COL_BLUE),
    ('2027年度末', 'DX1.0 主要施策 導入率70%', 'タイムスケジュール・データ収集・デジタル販促の定着確認', COL_BLUE),
    ('2028年度', 'ツルハHD データ基盤連携開始', 'ツルハHDデータ基盤刷新を受けた杏林堂側の本格連携開始', COL_GREEN),
    ('2030年度末', 'DX1.5 全社データ活用定着', '統合ダッシュボード・シフト予測・CRM施策の内製化完了', COL_GREEN),
    ('2031年以降', 'AI人財内製化・DX2.0始動', 'AI活用業務展開・CX推進部がAI人財組織としてグループ牽引', COL_PURPLE),
]
for i, (date, title, desc, col) in enumerate(ms_items):
    row = i // 2
    colx = i % 2
    x = 0.4 + colx * 6.5
    y = 0.85 + row * 2.1
    add_rect(slide, x, y, 6.1, 1.9, COL_WHITE)
    add_rect(slide, x, y, 0.08, 1.9, col)
    add_rect(slide, x, y, 6.1, 0.35, col)
    add_text(slide, date, x+0.15, y+0.05, 4, 0.27, font_size=10, bold=True, color=COL_WHITE)
    add_text(slide, title, x+0.15, y+0.42, 5.8, 0.35, font_size=12, bold=True, color=COL_TEXT)
    add_text(slide, desc, x+0.15, y+0.82, 5.7, 0.85, font_size=9.5, color=COL_TEXT)
print('スライド3完了')

# ========== スライド4: DX1.0 施策 ==========
slide = prs.slides.add_slide(blank_layout)
slide_bg(slide, COL_LIGHT)
add_rect(slide, 0, 0, 13.33, 0.6, COL_BLUE)
add_rect(slide, 0, 0.6, 13.33, 0.05, COL_ACCENT)
add_text(slide, 'DX 1.0 主要施策　〜2028年3月　業務自動化・デジタル基盤整備', 0.3, 0.1, 12, 0.45, font_size=18, bold=True, color=COL_WHITE)

kpis_10 = ['業務時間 20%削減', 'シフトシステム 全店100%', 'MS365活用 本部100%', 'データ自動取得率 60%']
kpi_x = 0.3
for kpi in kpis_10:
    add_rect(slide, kpi_x, 0.72, 3.1, 0.32, RGBColor(0xD6, 0xEA, 0xF8))
    add_text(slide, kpi, kpi_x+0.08, 0.74, 2.95, 0.28, font_size=8.5, bold=True, color=COL_BLUE)
    kpi_x += 3.26

施策_10 = [
    ('C', 'データ基盤', COL_C, [
        'データ収集ルール統一（店舗・販売・在庫）',
        'ツルハHD連携に向けた杏林堂側データ整備',
        'MS365 Power BI活用 データ可視化PoC',
    ]),
    ('A', 'シフト・人員管理', COL_A, [
        'タイムスケジュールシステム 全店展開',
        '応援依頼フローのシステム化・確立',
        '勤怠データとシフトデータの自動マッチング',
    ]),
    ('D', '顧客体験', COL_D, [
        '顧客データ収集基盤整備（購買・来店統合）',
        '販促施策のデジタル化（紙チラシ→デジタル）',
    ]),
    ('Ops', 'オペレーション', COL_GRAY, [
        '電子棚札の段階的展開',
        '発注自動化の検討・PoC',
        'Teams/Formsによる社内情報共有改善',
    ]),
]
for idx, (badge, theme_name, col, items) in enumerate(施策_10):
    x = col_positions[idx]
    w = 3.1
    y_top = 1.2
    add_rect(slide, x, y_top, w, 0.42, col)
    add_text(slide, badge, x+0.08, y_top+0.05, 0.5, 0.3, font_size=13, bold=True, color=COL_WHITE)
    add_text(slide, theme_name, x+0.55, y_top+0.08, w-0.65, 0.3, font_size=12, bold=True, color=COL_WHITE)
    add_rect(slide, x, y_top+0.42, w, 5.4, COL_WHITE)
    for i, item in enumerate(items):
        iy = y_top + 0.55 + i * 0.85
        add_rect(slide, x+0.12, iy, 0.15, 0.15, col)
        add_text(slide, item, x+0.35, iy-0.05, w-0.45, 0.7, font_size=9.5, color=COL_TEXT)

add_rect(slide, 0.25, 7.05, 12.85, 0.35, RGBColor(0xFD, 0xF2, 0xE9))
add_text(slide, '留意事項：MS365店舗展開（2027/3）完了後の活用定着・トレーニングコストに注意。OBIC API仕様確認が必要（未解決）', 0.4, 7.08, 12.5, 0.3, font_size=8.5, color=RGBColor(0xE6, 0x7E, 0x22))
print('スライド4完了')

# ========== スライド5: DX1.5 施策 ==========
slide = prs.slides.add_slide(blank_layout)
slide_bg(slide, COL_LIGHT)
add_rect(slide, 0, 0, 13.33, 0.6, COL_GREEN)
add_rect(slide, 0, 0.6, 13.33, 0.05, COL_ACCENT)
add_text(slide, 'DX 1.5 主要施策　2028〜2031年　データ活用・分析による意思決定支援', 0.3, 0.1, 12, 0.45, font_size=18, bold=True, color=COL_WHITE)

kpis_15 = ['データドリブン意思決定 70%', 'シフト予測精度 ±10%以内', '顧客分析活用 全店80%', 'レポート自動化 60%']
kpi_x = 0.3
for kpi in kpis_15:
    add_rect(slide, kpi_x, 0.72, 3.1, 0.32, RGBColor(0xD5, 0xF5, 0xE3))
    add_text(slide, kpi, kpi_x+0.08, 0.74, 2.95, 0.28, font_size=8.5, bold=True, color=COL_GREEN)
    kpi_x += 3.26

施策_15 = [
    ('C', 'データ基盤', COL_C, [
        'ツルハHD DB連携（販売・在庫・会員統合）',
        '店舗別・地域別 統合ダッシュボード構築',
        'データガバナンス・品質管理体制確立',
        'Power BI 自動レポーティング全社展開',
    ]),
    ('A', 'シフト・人員管理', COL_A, [
        '来客数・売上データによるシフト需要予測',
        '繁忙予測×スキルマッチングで応援依頼最適化',
        '人材育成データ（スキルマップ）のデジタル化',
    ]),
    ('D', '顧客体験', COL_D, [
        'CRMによる顧客セグメント分析',
        '購買データに基づく販促パーソナライズ化',
        'アプリ・ECとの顧客接点統合',
    ]),
    ('C+', 'データ連携基盤', COL_C, [
        'データ品質KPI設定・監視自動化',
        'BI活用トレーニング・全社展開',
        'データカタログ・メタデータ管理導入',
    ]),
]
for idx, (badge, theme_name, col, items) in enumerate(施策_15):
    x = col_positions[idx]
    w = 3.1
    y_top = 1.2
    add_rect(slide, x, y_top, w, 0.42, col)
    add_text(slide, badge, x+0.08, y_top+0.05, 0.5, 0.3, font_size=13, bold=True, color=COL_WHITE)
    add_text(slide, theme_name, x+0.55, y_top+0.08, w-0.65, 0.3, font_size=12, bold=True, color=COL_WHITE)
    add_rect(slide, x, y_top+0.42, w, 5.4, COL_WHITE)
    for i, item in enumerate(items):
        iy = y_top + 0.55 + i * 0.85
        add_rect(slide, x+0.12, iy, 0.15, 0.15, col)
        add_text(slide, item, x+0.35, iy-0.05, w-0.45, 0.7, font_size=9.5, color=COL_TEXT)

add_rect(slide, 0.25, 7.05, 12.85, 0.35, RGBColor(0xEA, 0xF7, 0xEE))
add_text(slide, '前提条件：ツルハHDデータ基盤刷新のスケジュール確定が中期計画の精度に影響。DX1.0のデータ整備完了が前提。', 0.4, 7.08, 12.5, 0.3, font_size=8.5, color=COL_GREEN)
print('スライド5完了')

# ========== スライド6: DX2.0 施策 ==========
slide = prs.slides.add_slide(blank_layout)
slide_bg(slide, COL_LIGHT)
add_rect(slide, 0, 0, 13.33, 0.6, COL_PURPLE)
add_rect(slide, 0, 0.6, 13.33, 0.05, COL_ACCENT)
add_text(slide, 'DX 2.0 主要施策　2031年以降　AI活用・自律最適化・顧客体験革新', 0.3, 0.1, 12, 0.45, font_size=18, bold=True, color=COL_WHITE)

kpis_20 = ['AI支援業務 50%以上', '顧客満足度 業界上位水準', 'シフト自動生成 90%以上', 'AIリテラシー保有 30%']
kpi_x = 0.3
for kpi in kpis_20:
    add_rect(slide, kpi_x, 0.72, 3.1, 0.32, RGBColor(0xE8, 0xDA, 0xEF))
    add_text(slide, kpi, kpi_x+0.08, 0.74, 2.95, 0.28, font_size=8.5, bold=True, color=COL_PURPLE)
    kpi_x += 3.26

施策_20 = [
    ('C', 'データ基盤', COL_C, [
        '全社データプラットフォーム 高度活用（リアルタイム分析）',
        '外部データ（人口動態・競合・天候）統合分析',
        'データ収益化・新規価値創造の検討',
    ]),
    ('A', 'シフト・人員管理', COL_A, [
        '人員需要自動予測・最適配置AI 全店導入',
        'スキル・資格・希望を考慮した自動シフト生成',
        '採用計画へのデータ活用（需要予測×人材）',
    ]),
    ('D', '顧客体験', COL_D, [
        '顧客個別最適化（レコメンド・健康アドバイス）',
        'AIを活用した調剤・健康相談支援',
        'オムニチャネルのシームレス体験実現',
    ]),
    ('E', '組織・人材', COL_E, [
        'AI人財の組織内自立（内製化・育成完了）',
        'CX推進部モデルの他部門・グループへの横展開',
        'AI倫理・ガバナンス体制の確立',
    ]),
]
for idx, (badge, theme_name, col, items) in enumerate(施策_20):
    x = col_positions[idx]
    w = 3.1
    y_top = 1.2
    add_rect(slide, x, y_top, w, 0.42, col)
    add_text(slide, badge, x+0.08, y_top+0.05, 0.5, 0.3, font_size=13, bold=True, color=COL_WHITE)
    add_text(slide, theme_name, x+0.55, y_top+0.08, w-0.65, 0.3, font_size=12, bold=True, color=COL_WHITE)
    add_rect(slide, x, y_top+0.42, w, 5.4, COL_WHITE)
    for i, item in enumerate(items):
        iy = y_top + 0.55 + i * 0.85
        add_rect(slide, x+0.12, iy, 0.15, 0.15, col)
        add_text(slide, item, x+0.35, iy-0.05, w-0.45, 0.7, font_size=9.5, color=COL_TEXT)

add_rect(slide, 0.25, 7.05, 12.85, 0.35, RGBColor(0xF0, 0xE6, 0xF6))
add_text(slide, 'DX2.0はDX1.5の成果・人財育成の達成度を踏まえ、開始タイミング・優先施策を見直す。', 0.4, 7.08, 12.5, 0.3, font_size=8.5, color=COL_PURPLE)
print('スライド6完了')

# ========== スライド7: KPI一覧 ==========
slide = prs.slides.add_slide(blank_layout)
slide_bg(slide, COL_LIGHT)
add_rect(slide, 0, 0, 13.33, 0.6, COL_DARK)
add_rect(slide, 0, 0.6, 13.33, 0.05, COL_ACCENT)
add_text(slide, 'フェーズ別 KPI目標一覧', 0.3, 0.1, 10, 0.45, font_size=20, bold=True, color=COL_WHITE)

headers = ['KPI指標', 'DX 1.0（〜2028/3）', 'DX 1.5（2028〜2031）', 'DX 2.0（2031〜）']
col_widths = [3.4, 3.1, 3.1, 3.1]
col_xs = [0.35, 3.85, 7.05, 10.25]
header_colors = [COL_DARK, COL_BLUE, COL_GREEN, COL_PURPLE]
header_y = 0.8

for i, (hdr, w, x, hcol) in enumerate(zip(headers, col_widths, col_xs, header_colors)):
    add_rect(slide, x, header_y, w, 0.42, hcol)
    add_text(slide, hdr, x+0.08, header_y+0.07, w-0.16, 0.3, font_size=11, bold=True, color=COL_WHITE, align=PP_ALIGN.CENTER)

kpi_rows = [
    ('業務時間削減率', '主要業務で20%削減', '主要業務で35%削減', '主要業務で50%削減'),
    ('シフト管理', '全店システム化100%', '予測精度 ±10%以内', 'AI自動生成率 90%以上'),
    ('データ活用', 'KPI自動取得率 60%', 'データドリブン決定 70%', 'AI支援業務 50%以上'),
    ('顧客分析', '収集基盤整備完了', '全店販促データ活用 80%', '個別最適化・NPS業界上位'),
    ('レポート自動化', '主要KPIのBI可視化', '定例報告 60%自動生成', '全社レポート 80%自動化'),
    ('人材・組織', 'MS365活用 本部100%', 'BIトレーニング全社展開', 'AIリテラシー 本部30%以上'),
]
for r, row_data in enumerate(kpi_rows):
    y = header_y + 0.42 + r * 0.95
    bg_even = RGBColor(0xFF, 0xFF, 0xFF)
    bg_odd = RGBColor(0xF8, 0xF9, 0xFA)
    for i, (text, w, x) in enumerate(zip(row_data, col_widths, col_xs)):
        bg = RGBColor(0xEB, 0xF5, 0xFB) if i == 0 else (bg_even if r % 2 == 0 else bg_odd)
        add_rect(slide, x, y, w, 0.93, bg)
        col_text = COL_BLUE if i == 0 else COL_TEXT
        add_text(slide, text, x+0.08, y+0.1, w-0.16, 0.75, font_size=9.5, color=col_text, bold=(i == 0))
print('スライド7完了')

# ========== スライド8: 推進体制・留意事項 ==========
slide = prs.slides.add_slide(blank_layout)
slide_bg(slide, COL_LIGHT)
add_rect(slide, 0, 0, 13.33, 0.6, COL_DARK)
add_rect(slide, 0, 0.6, 13.33, 0.05, COL_ACCENT)
add_text(slide, '推進体制・留意事項・前提条件', 0.3, 0.1, 10, 0.45, font_size=20, bold=True, color=COL_WHITE)

add_rect(slide, 0.3, 0.8, 12.7, 0.35, COL_DARK)
add_text(slide, '推進体制', 0.5, 0.83, 5, 0.28, font_size=12, bold=True, color=COL_WHITE)
taisei_items = [
    ('推進主体', 'CX推進部（各テーマリード担当を設定）'),
    ('意思決定', '四半期ごとにロードマップの進捗確認・見直し'),
    ('報告体系', '経営層への定期報告（半期）＋チーム内週次進捗管理'),
]
for i, (label, content) in enumerate(taisei_items):
    y = 1.2 + i * 0.55
    add_rect(slide, 0.3, y, 2.8, 0.48, RGBColor(0xEB, 0xF5, 0xFB))
    add_text(slide, label, 0.45, y+0.08, 2.5, 0.3, font_size=10, bold=True, color=COL_BLUE)
    add_rect(slide, 3.15, y, 9.85, 0.48, COL_WHITE)
    add_text(slide, content, 3.3, y+0.08, 9.5, 0.3, font_size=10, color=COL_TEXT)

add_rect(slide, 0.3, 2.95, 12.7, 0.35, COL_GREEN)
add_text(slide, '前提条件', 0.5, 2.98, 5, 0.28, font_size=12, bold=True, color=COL_WHITE)
prereqs = [
    ('ツルハHDデータ基盤刷新', 'スケジュール確定が中期計画（DX1.5）の精度に直接影響'),
    ('MS365店舗展開', '2027年3月完了が前提。展開後の活用定着・トレーニングが必要'),
    ('OBIC連携', 'シフト管理アプリ選定においてOBIC API仕様の確認が必要（未解決）'),
]
for i, (label, content) in enumerate(prereqs):
    y = 3.35 + i * 0.55
    add_rect(slide, 0.3, y, 2.8, 0.48, RGBColor(0xEA, 0xF7, 0xEE))
    add_text(slide, label, 0.45, y+0.08, 2.5, 0.3, font_size=10, bold=True, color=COL_GREEN)
    add_rect(slide, 3.15, y, 9.85, 0.48, COL_WHITE)
    add_text(slide, content, 3.3, y+0.08, 9.5, 0.3, font_size=10, color=COL_TEXT)

add_rect(slide, 0.3, 5.12, 12.7, 0.35, RGBColor(0xC0, 0x39, 0x2B))
add_text(slide, 'リスクと対応', 0.5, 5.15, 5, 0.28, font_size=12, bold=True, color=COL_WHITE)
risks = [
    ('変化への抵抗', '店舗スタッフのデジタルリテラシー格差。段階的導入と研修計画が必要'),
    ('データ品質', 'DX1.0段階でのデータ収集精度が後続フェーズに影響。早期に品質基準設定'),
]
for i, (label, content) in enumerate(risks):
    y = 5.52 + i * 0.55
    add_rect(slide, 0.3, y, 2.8, 0.48, RGBColor(0xFD, 0xED, 0xEC))
    add_text(slide, label, 0.45, y+0.08, 2.5, 0.3, font_size=10, bold=True, color=RGBColor(0xC0, 0x39, 0x2B))
    add_rect(slide, 3.15, y, 9.85, 0.48, COL_WHITE)
    add_text(slide, content, 3.3, y+0.08, 9.5, 0.3, font_size=10, color=COL_TEXT)

add_rect(slide, 0, 7.3, 13.33, 0.2, COL_DARK)
add_text(slide, '杏林堂薬局 CX推進部 DXロードマップ | 2026年4月版', 0.3, 7.31, 10, 0.18, font_size=8, color=COL_GRAY)
print('スライド8完了')

# 保存
out_dir = os.path.join(os.getcwd(), 'kyorindo-cx', '07_deliverables')
os.makedirs(out_dir, exist_ok=True)
out_path = os.path.join(out_dir, 'CX_DX_Roadmap_20260405.pptx')
prs.save(out_path)
print(f'保存完了: {out_path}')
print(f'ファイルサイズ: {os.path.getsize(out_path):,} bytes')
