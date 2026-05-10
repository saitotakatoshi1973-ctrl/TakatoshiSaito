"""
gemini_wiki_generator.py
========================
Excel / PPTX / PDF ファイルを Gemini 2.5 Flash API で wiki Markdown に変換するスクリプト。

【スタンドアロンモード】ファイルを指定して直接wiki生成・保存
    python gemini_wiki_generator.py <ファイルパス> [--dest <保存先>] [--dry-run]

【分類モード（--analyze-only）】batch-inbox.md から呼び出し、分類結果JSONのみ stdout に出力
    python gemini_wiki_generator.py <ファイルパス> --analyze-only
        [--schema-path <SCHEMA.md>] [--hints-path <classification-hints.md>] [--emit-usage]

【生成モード（--generate）】分類結果を受け取り、本文Markdownのみ stdout に出力
    python gemini_wiki_generator.py <ファイルパス> --generate
        [--analysis-json <分類結果JSON>] [--analysis-json-stdin] [--emit-usage]

【本番保存モード（--save-wiki）】分類・本文生成・wiki保存・processed-sources更新を行う
    python gemini_wiki_generator.py <ファイルパス> --save-wiki [--emit-usage]

【互換モード（--body-only）】write-wiki.md から呼び出し、本文のみ stdout に出力
    python gemini_wiki_generator.py <ファイルパス> --body-only
        [--wiki-type strategy] [--title "タイトル"]
        [--scope kyorindo] [--domain cx] [--tags "CX,戦略"]
        [--converted-text "事前抽出テキスト"]

例:
    python gemini_wiki_generator.py "C:/Users/.../資料.xlsx"
    python gemini_wiki_generator.py "C:/Users/.../資料.pptx" --dest kyorindo/it-systems/
    python gemini_wiki_generator.py "C:/Users/.../資料.pdf" --dry-run
    python gemini_wiki_generator.py "C:/Users/.../資料.pptx" --body-only --wiki-type strategy --title "CX戦略"
"""

import argparse
import hashlib
import json
import os
import re
import sys
import textwrap
from datetime import date
from pathlib import Path
from tempfile import NamedTemporaryFile

import yaml
from dotenv import load_dotenv

# ── 定数 ────────────────────────────────────────────────
PERSONAL_ROOT  = Path(r"C:\Users\takatoshi-saito\OneDrive\00personal")
KB_ROOT        = PERSONAL_ROOT / "KnowledgeBase"
PROCESSED_PATH = KB_ROOT / "_system" / "processed-sources.yaml"
DOTENV_PATH    = PERSONAL_ROOT / "ClaudeCodeFolder" / ".env"
SCHEMA_PATH    = KB_ROOT / "_system" / "SCHEMA.md"
HINTS_PATH     = KB_ROOT / "_system" / "learning" / "classification-hints.md"

TODAY = date.today().isoformat()
VALID_WIKI_TYPES = [
    "strategy",
    "organization",
    "issue",
    "progress",
    "reference",
    "task",
    "log",
    "benchmark",
    "research",
    "specification",
    "report",
    "manual",
    "meeting_notes",
]


# ── .env 読み込み ────────────────────────────────────────
load_dotenv(DOTENV_PATH)
GEMINI_API_KEY = os.getenv("GEMINI_API_KEY")
if not GEMINI_API_KEY:
    print("❌ GEMINI_API_KEY が設定されていません。.env ファイルを確認してください。")
    sys.exit(1)


# ── Gemini クライアント初期化 ────────────────────────────
from google import genai
CLIENT = genai.Client(api_key=GEMINI_API_KEY)
MODEL_NAME = "gemini-2.5-flash"


# ════════════════════════════════════════════════════════
# テキスト抽出
# ════════════════════════════════════════════════════════

def extract_text(file_path: Path) -> str:
    """拡張子に応じてテキストを抽出する"""
    ext = file_path.suffix.lower()

    if ext == ".xlsx":
        return _extract_excel(file_path)
    elif ext == ".pptx":
        return _extract_pptx(file_path)
    elif ext == ".pdf":
        return _extract_pdf(file_path)
    elif ext in (".md", ".txt"):
        return file_path.read_text(encoding="utf-8", errors="ignore")
    else:
        raise ValueError(f"未対応の形式: {ext}")


def _extract_excel(file_path: Path) -> str:
    """Excel ファイルからテキストを抽出する"""
    import openpyxl
    wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
    lines = []
    for ws in wb.worksheets:
        lines.append(f"## シート: {ws.title}")
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            # 空行はスキップ
            if any(c.strip() for c in cells):
                lines.append("\t".join(cells))
    wb.close()
    return "\n".join(lines)


def _extract_pptx(file_path: Path) -> str:
    """PPTX ファイルからテキストを抽出する"""
    from pptx import Presentation
    prs = Presentation(file_path)
    lines = []
    for i, slide in enumerate(prs.slides, 1):
        slide_texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        slide_texts.append(text)
        if slide_texts:
            lines.append(f"## スライド {i}")
            lines.extend(slide_texts)
    return "\n".join(lines)


def _extract_pdf(file_path: Path) -> str:
    """PDF ファイルからテキストを抽出する"""
    import fitz  # PyMuPDF
    doc = fitz.open(file_path)
    lines = []
    for i, page in enumerate(doc, 1):
        text = page.get_text().strip()
        if text:
            lines.append(f"## ページ {i}")
            lines.append(text)
    doc.close()
    return "\n".join(lines)


# ════════════════════════════════════════════════════════
# Gemini API でwiki生成
# ════════════════════════════════════════════════════════

def _usage_dict(response, *, mode: str, input_text: str, output_text: str) -> dict:
    """
    Gemini API の usage_metadata を、ログやJSON出力に載せやすい辞書へ正規化する。

    Codex/Claude 側コストを下げるには、重い分類・本文生成をこの Python 側へ寄せる必要がある。
    その効果を確認するため、Gemini 側で実際に使ったトークン数をできる限り記録する。
    SDKやモデルによって usage_metadata の属性名が変わる可能性があるため、getattr で安全に読む。
    """
    usage = getattr(response, "usage_metadata", None)
    return {
        "mode": mode,
        "model": MODEL_NAME,
        "input_chars": len(input_text),
        "output_chars": len(output_text),
        "prompt_token_count": getattr(usage, "prompt_token_count", None) if usage else None,
        "candidates_token_count": getattr(usage, "candidates_token_count", None) if usage else None,
        "total_token_count": getattr(usage, "total_token_count", None) if usage else None,
    }


def _generate_content_with_usage(prompt: str, *, mode: str) -> tuple[str, dict]:
    """
    Gemini 呼び出しを1か所に集約する。

    各モードが個別に CLIENT.models.generate_content() を呼ぶと、usageログの形式が揺れやすい。
    ここで本文と usage を同時に返すことで、analyze/generate/dry-run のコスト比較を同じ形式で行う。
    """
    response = CLIENT.models.generate_content(
        model=MODEL_NAME,
        contents=prompt,
    )
    text = response.text or ""
    return text, _usage_dict(response, mode=mode, input_text=prompt, output_text=text)


def _print_usage(usage: dict) -> None:
    """
    usage は stderr に出す。

    stdout は --analyze-only のJSONや --body-only/--generate の本文を機械的に受け渡すために使う。
    ここへログを混ぜると batch-inbox.md 側のパースを壊すので、計測ログは stderr に逃がす。
    """
    print(json.dumps({"usage": usage}, ensure_ascii=False), file=sys.stderr)


def _extract_category_path_entries(schema_text: str) -> list[dict]:
    """
    SCHEMA.md 内の fenced YAML から category_paths を抽出する。

    SCHEMA.md の3.x各カテゴリには、人間にも読みやすく機械にも安定して読めるよう、
    `category_paths:` YAMLブロックを置く運用にする。
    Python側ではこのYAMLを最優先し、allow_wiki=true の path だけを保存先候補として使う。
    allow_wiki=false の `_system/` などは構造としては残すが、Geminiの保存先候補には渡さない。
    """
    entries: list[dict] = []
    for match in re.finditer(r"```yaml\s*(.*?)\s*```", schema_text, flags=re.DOTALL):
        block = match.group(1)
        if "category_paths:" not in block:
            continue
        try:
            data = yaml.safe_load(block) or {}
        except yaml.YAMLError:
            continue
        for item in data.get("category_paths", []) or []:
            if not isinstance(item, dict):
                continue
            path = str(item.get("path", "")).strip()
            if not path:
                continue
            allow_wiki = item.get("allow_wiki", True)
            entries.append({
                "path": path.strip("/") + "/",
                "description": str(item.get("description", "")).strip(),
                "allow_wiki": bool(allow_wiki),
            })
    return entries


def _extract_schema_destinations(schema_text: str) -> list[str]:
    """
    SCHEMA.md から Gemini に選ばせてよい保存先候補を抽出する。

    Gemini に自由にパスを生成させると、存在しないカテゴリや表記揺れを作る危険がある。
    そのため分類プロンプトには「この候補から選ぶ」制約を渡す。
    SCHEMA.md には2種類の表現がある。
    1. `kyorindo/it-systems/infrastructure/` のような明示パス
    2. Markdownのコードブロック内ツリー:
       `├── it-systems/`、`│   ├── infrastructure/`

    最初の実装は英数字の明示パスだけを拾っていたため、日本語説明が中心のSCHEMAでは
    保存先候補が不足し、Geminiが destination="" を返しやすかった。
    ここでは明示パスに加えて、ツリー構造から親子関係を復元する。
    """
    yaml_entries = _extract_category_path_entries(schema_text)
    yaml_candidates = sorted({
        entry["path"]
        for entry in yaml_entries
        if entry.get("allow_wiki", True)
    })
    if yaml_candidates:
        return yaml_candidates

    candidates: set[str] = set()
    top_levels = {"kyorindo", "tsuruha-hd", "retail", "ai-dx", "vendor", "research"}

    def add_candidate(path: str) -> None:
        """
        SCHEMAから拾ったパス候補を正規化して登録する。

        SCHEMA.mdには説明文・ファイル名・例示が混ざるため、何でも候補化すると
        Geminiが存在しない保存先を選びやすくなる。ここでは wiki のトップカテゴリから
        始まるパスだけを候補にし、`.md` や `.yaml` の運用ファイルは除外する。
        """
        normalized = path.strip().strip("/").replace("\\", "/")
        if not normalized:
            return
        if any(part.endswith((".md", ".yaml", ".yml")) for part in normalized.split("/")):
            return
        if "（" in normalized or "）" in normalized:
            return
        if normalized.split("/", 1)[0] not in top_levels:
            return
        candidates.add(normalized + "/")

    for match in re.finditer(r"([^\s`]+(?:/[^\s`]+)+/?)", schema_text):
        add_candidate(match.group(1))

    stack: list[str] = []
    in_tree = False
    for raw_line in schema_text.splitlines():
        line = raw_line.rstrip()
        if line.strip() == "```":
            in_tree = not in_tree
            stack = []
            continue
        if not in_tree or "──" not in line:
            continue

        match = re.search(r"(?:├──|└──)\s+([^#]+)", line)
        if not match:
            continue
        name = match.group(1).strip()
        name = name.split("#", 1)[0].strip()
        name = name.rstrip("/")
        if not name or name.startswith("..."):
            continue

        # ツリーのインデントは概ね4文字単位（"│   " や空白）で深くなる。
        # 多少の表記揺れを許容するため、罫線より前の文字数から階層を推定する。
        prefix = line[: match.start()]
        depth = max(0, len(prefix) // 4)
        stack = stack[:depth]
        stack.append(name)

        add_candidate("/".join(stack))
    return sorted(candidates)


def _format_destination_prompt(schema_text: str, destinations: list[str]) -> str:
    """
    Geminiへ渡す保存先候補テキストを作る。

    YAMLの category_paths がある場合は description も渡す。
    pathだけよりも `path: description` の方が分類判断に必要な意味情報が増えるため、
    `複合機` → `kyorindo/it-systems/infrastructure/` のような判断が安定しやすい。
    """
    entries = _extract_category_path_entries(schema_text)
    if entries:
        lines = []
        for entry in entries:
            if not entry.get("allow_wiki", True):
                continue
            desc = entry.get("description", "")
            if desc:
                lines.append(f"- {entry['path']}: {desc}")
            else:
                lines.append(f"- {entry['path']}")
        return "\n".join(lines) or "（保存先候補なし）"
    return "\n".join(f"- {p}" for p in destinations) or "（保存先候補なし）"


def _destination_matches_schema(destination: str, candidates: list[str]) -> bool:
    """
    Geminiが返したdestinationがSCHEMA候補に沿っているかを判定する。

    `vendor/{vendor-name}/` のようなテンプレートカテゴリは、実運用では
    `vendor/xerox/` のような具体名に展開される。完全一致だけで判定すると
    正しい分類まで review 扱いになるため、`{...}` を1階層ワイルドカードとして扱う。
    """
    normalized = destination.strip().strip("/") + "/"
    if normalized in candidates:
        return True
    for candidate in candidates:
        pattern = re.escape(candidate.strip("/"))
        pattern = re.sub(r"\\\{[^/]+\\\}", r"[^/]+", pattern)
        if re.fullmatch(pattern + r"/?", normalized.strip("/")):
            return True
    return False


def _load_text_if_exists(path: Path) -> str:
    """設定ファイルが存在しない環境でも dry-run できるよう、安全に読み込む。"""
    if not path.exists():
        return ""
    return path.read_text(encoding="utf-8")


def _json_from_model_text(text: str) -> dict:
    """
    Gemini から返った分類JSONを辞書化する。

    モデルが ```json ... ``` で囲むことがあるため、まずコードフェンスを剥がす。
    それでもJSONでない場合は明示的にエラーにして、batch-inbox.md 側でフォールバック判断できるようにする。
    """
    cleaned = text.strip()
    if cleaned.startswith("```"):
        cleaned = re.sub(r"^```(?:json)?\s*", "", cleaned)
        cleaned = re.sub(r"\s*```$", "", cleaned)
    return json.loads(cleaned)


def _load_analysis_json(value: str) -> dict:
    """
    --generate の --analysis-json を読み込む。

    batch-inbox.md からは、分類結果JSONをそのまま渡す場合と、
    一時ファイルに保存したJSONパスを渡す場合の両方があり得る。
    JSON文字列を Path として解釈すると Windows では長すぎるパス扱いで失敗することがあるため、
    先頭が "{" または "[" なら即JSONとして読む。
    """
    stripped = value.strip()
    if not stripped:
        return {}
    if stripped.startswith("{") or stripped.startswith("["):
        return json.loads(stripped)
    candidate = Path(stripped)
    if candidate.exists():
        return json.loads(candidate.read_text(encoding="utf-8"))
    return json.loads(stripped)


def _load_analysis_from_args(args) -> dict:
    """
    --generate 用の分類結果JSONを読み込む。

    優先順位:
    1. --analysis-json-stdin
       PowerShellやバッチ処理では、複数行JSONをコマンドライン引数にすると
       空白や改行で分割されて argparse が壊れやすい。
       標準入力ならJSONをそのまま渡せるため、バッチ実装ではこの方式を推奨する。

    2. --analysis-json
       既存互換として残す。短いJSON文字列またはJSONファイルパスを受け付ける。

    3. 空辞書
       明示分類なしで --wiki-type 等の個別CLI引数へフォールバックする。
    """
    if getattr(args, "analysis_json_stdin", False):
        raw = sys.stdin.read()
        if not raw.strip():
            return {}
        return json.loads(raw)
    if args.analysis_json:
        return _load_analysis_json(args.analysis_json)
    return {}


ANALYZE_PROMPT_TEMPLATE = textwrap.dedent("""\
    あなたは社内ナレッジベースの分類担当です。
    以下の文書を読み、wiki化する際の分類情報をJSONだけで返してください。

    ## 重要な制約

    - 出力はJSONオブジェクトのみ。Markdownや説明文を前後に付けない。
    - destination は必ず「保存先候補」の中から1つだけ選ぶ。
    - 保存先候補に適切なものがない場合は、destination を空文字にし、needs_review を true にする。
    - wiki_type は許可リストから1つだけ選ぶ。
    - 不確実な場合は confidence_score を低くし、needs_review を true にする。
    - classification-hints.md に合うルールがある場合は、それを最優先する。
    - 元資料にない事実を推測して分類理由に書かない。

    ## JSON形式

    {{
      "status": "success",
      "title": "文書タイトル",
      "wiki_type": "reference",
      "destination": "kyorindo/it-systems/example/",
      "tags": ["タグ1", "タグ2"],
      "confidence_score": 0,
      "classification_method": "gemini_schema_hints",
      "reason": "分類理由",
      "needs_review": true
    }}

    ## wiki_type 許可リスト

    {wiki_types}

    ## 保存先候補

    {destinations}

    ## classification-hints.md

    {hints}

    ## 文書情報

    ファイル名: {file_name}
    フォルダパス: {folder_path}

    ## 文書テキスト

    {content}
""")

# ── wiki_type 別テンプレート ─────────────────────────────
WIKI_TYPE_TEMPLATES = {
    "strategy": """\
## 背景・目的
（元資料から読み取れる背景や目的を2〜4文で記述）

## 主要内容
（ロードマップ・施策・方針を箇条書きまたはフェーズ別に整理）

## 重要ポイント
（意思決定に影響する数値・KPI・優先事項を箇条書き）

## 関連リンク
（分類先フォルダ内の関連ファイルがあればパスを記述、なければ空）""",

    "organization": """\
## 概要
（組織・企業の基本情報）

## 組織構成
（部署・チーム・人員構成を箇条書きまたは表で整理）

## 役割・担当
（主要ポストと担当業務）

## 変更履歴
（直近の組織変更があれば記録）""",

    "issue": """\
## 課題一覧

| # | 課題 | 優先度 | 担当 |
|---|------|--------|------|
| 1 | ...  | 高     | ...  |

## 背景・原因
（課題の発生背景や根本原因）

## 対応方針
（解決に向けた方向性・アプローチ）""",

    "progress": """\
## 日時・参加者
- 日時: YYYY-MM-DD HH:MM
- 参加者: （氏名・部署）

## 議題
（議題を箇条書き）

## 決定事項
（決定内容を箇条書き）

## 次回アクション

| 担当 | 内容 | 期限 |
|------|------|------|""",

    "reference": """\
## 出典情報
- 媒体/著者:
- 公開日:

## 要点サマリー
（3〜5文で主要な主張・結論を要約）

## 詳細メモ
（重要な数値・固有名詞・ファクトを箇条書き）

## 自社への示唆
（杏林堂やプロジェクトに関連する示唆があれば記述）""",

    "task": """\
## タスク一覧

| # | タスク | 優先度 | 担当 | 期限 | 状態 |
|---|--------|--------|------|------|------|

## 背景・目的
（このタスクリストの目的）""",

    "log": """\
## 日時
YYYY-MM-DD

## 内容
（議論や作業の概要）

## 判断理由
（なぜその判断をしたか）

## 次のステップ
（続きのアクション）""",

    "benchmark": """\
## 比較軸
（何を比較するかの定義）

## 比較テーブル

| 企業/項目 | 指標1 | 指標2 | 指標3 |
|----------|-------|-------|-------|

## 考察
（比較から読み取れる示唆・杏林堂の立ち位置）""",

    "research": """\
## 出典
- 主催/著者:
- 日時/公開日:
- 形式: セミナー | レポート | 論文

## 概要
（2〜3文でテーマと主旨を説明）

## 重要な知見
（箇条書きで重要な発見・主張を列挙）

## 自社への示唆
（実務・プロジェクトへの応用可能性）""",
}

# ── スタンドアロン用プロンプト（Front-matter込み） ───────────
WIKI_PROMPT_TEMPLATE = textwrap.dedent("""\
    あなたは社内ナレッジベースのwikiライターです。
    以下のビジネス文書のテキストを読み、社内wiki用のMarkdownファイルを生成してください。

    ## 出力形式

    必ず以下のYAML front-matterから始めてください：

    ---
    title: "（文書のタイトル）"
    wiki_type: "（meeting_notes / reference / strategy / specification / report / manual のいずれか）"
    source: "{source_path}"
    source_date: "（文書の作成日・更新日。YYYYMMDD形式。不明なら空文字）"
    created: "{today}"
    status: active
    tags: [（関連タグをカンマ区切りで）]
    ---

    ## 本文の書き方

    - 文書の内容を日本語で整理してください
    - 重要な数値・日付・固有名詞は必ず含めてください
    - テーブル形式で整理できる情報はMarkdownテーブルにしてください
    - 見出し（##）で構造化してください
    - 最後に「## 関連情報」セクションを追加してください（関連が推測できる場合）

    ## 文書テキスト

    ファイル名: {file_name}
    フォルダパス: {folder_path}

    {content}
""")

# ── スキル統合用プロンプト（本文のみ・Front-matterなし） ───────
WIKI_BODY_PROMPT_TEMPLATE = textwrap.dedent("""\
    あなたは社内ナレッジベースのwikiライターです。
    以下のビジネス文書のテキストを読み、wiki記事の本文を日本語で執筆してください。

    ## 分類情報（システムが事前に決定済み）

    wiki_type : {wiki_type}
    タイトル  : {title}
    保存先    : {destination}
    スコープ  : {scope}
    ドメイン  : {domain}
    タグ      : {tags}

    ## 出力セクション構成（wiki_typeに対応するテンプレート）

    {template}

    ## 執筆ルール

    - Front-matter（---で囲まれたYAMLブロック）は出力しない
    - テンプレートのセクション構成を必ず守る
    - 元資料にない内容は推測して書かない
    - 元資料に情報が不足するセクションは「（情報なし）」と記載する
    - 固有名詞・数値・日付は元資料のものを正確に使う
    - 文体は「である調」で統一する
    - 分量: 300〜800文字（summaryモード時は150〜300文字）

    ## 文書テキスト

    ファイル名: {file_name}
    フォルダパス: {folder_path}

    {content}
""")


def analyze_document(
    file_path: Path,
    content: str,
    schema_path: Path = SCHEMA_PATH,
    hints_path: Path = HINTS_PATH,
) -> tuple[dict, dict]:
    """
    【分類モード】Geminiで保存先・wiki_type・タイトルを決める。

    これは analyze.md の低コスト代替として batch-inbox.md から使う想定の関数。
    analyze.md は分類精度の高いフォールバックとして残し、通常バッチではまずこの関数を使う。
    """
    schema_text = _load_text_if_exists(schema_path)
    hints_text = _load_text_if_exists(hints_path)
    destinations = _extract_schema_destinations(schema_text)
    destination_text = _format_destination_prompt(schema_text, destinations)
    prompt = ANALYZE_PROMPT_TEMPLATE.format(
        wiki_types=", ".join(VALID_WIKI_TYPES),
        destinations=destination_text,
        hints=hints_text[:12000] or "（なし）",
        file_name=file_path.name,
        folder_path=str(file_path.parent),
        content=content[:50000],
    )
    raw_text, usage = _generate_content_with_usage(prompt, mode="analyze")
    result = _json_from_model_text(raw_text)

    # Gemini が候補外パスを返した場合は、そのまま採用せずレビュー扱いにする。
    # ここで強制的に落とすより、JSONに状態を残す方がバッチ後レビューしやすい。
    destination = result.get("destination", "")
    if destination and destinations and not _destination_matches_schema(destination, destinations):
        result["needs_review"] = True
        result["destination_error"] = "destination_not_in_schema_candidates"

    if result.get("wiki_type") not in VALID_WIKI_TYPES:
        result["needs_review"] = True
        result["wiki_type_error"] = "wiki_type_not_allowed"

    result.setdefault("status", "success")
    result.setdefault("classification_method", "gemini_schema_hints")
    result["schema_candidate_count"] = len(destinations)
    result["schema_candidate_sample"] = destinations[:20]
    return result, usage


def generate_wiki(file_path: Path, content: str) -> tuple[str, dict]:
    """【スタンドアロン用】Front-matter込みのwiki Markdownを生成する"""
    source_path = to_personal_relative(file_path)
    prompt = WIKI_PROMPT_TEMPLATE.format(
        source_path=source_path,
        today=TODAY,
        file_name=file_path.name,
        folder_path=str(file_path.parent),
        content=content[:50000],
    )
    print("🤖 Gemini API に送信中...")
    return _generate_content_with_usage(prompt, mode="standalone")


def generate_wiki_body(
    file_path: Path,
    content: str,
    wiki_type: str = "reference",
    title: str = "",
    destination: str = "",
    scope: str = "",
    domain: str = "",
    tags: list[str] | None = None,
) -> tuple[str, dict]:
    """【スキル統合用】本文のみ（Front-matterなし）を生成して返す。
    analyze.md の分類結果を受け取り、wiki_type に対応するテンプレートで執筆する。
    write-wiki.md から --body-only モードで呼び出される。
    """
    template = WIKI_TYPE_TEMPLATES.get(wiki_type, WIKI_TYPE_TEMPLATES["reference"])
    prompt = WIKI_BODY_PROMPT_TEMPLATE.format(
        wiki_type=wiki_type,
        title=title or file_path.stem,
        destination=destination,
        scope=scope,
        domain=domain,
        tags=", ".join(tags or []),
        template=template,
        file_name=file_path.name,
        folder_path=str(file_path.parent),
        content=content[:50000],
    )
    return _generate_content_with_usage(prompt, mode="generate")


# ════════════════════════════════════════════════════════
# wiki ファイル保存
# ════════════════════════════════════════════════════════

def suggest_wiki_filename(file_path: Path) -> str:
    """ファイル名からwikiファイル名を推定する（英小文字・アンダースコア）"""
    stem = file_path.stem
    # 日付パターンを保持しつつ、日本語・記号を除去
    stem = re.sub(r"[^\w\s\-]", "_", stem, flags=re.UNICODE)
    stem = re.sub(r"\s+", "_", stem)
    stem = re.sub(r"_+", "_", stem).strip("_").lower()
    # ASCII以外を除去
    stem = re.sub(r"[^\x00-\x7F]", "", stem).strip("_")
    if not stem:
        stem = "wiki_" + hashlib.md5(file_path.name.encode()).hexdigest()[:8]
    return stem + ".md"


def suggest_wiki_filename_from_title(title: str, source_path: Path) -> str:
    """
    wikiタイトルから保存ファイル名を作る。

    既存の suggest_wiki_filename() はASCII寄りで、日本語タイトルが消えやすい。
    KnowledgeBaseには日本語ファイル名が既に使われているため、本番保存では日本語を残す。
    Windowsで使えない記号だけを安全なアンダースコアへ置換し、日付を付けて衝突しにくくする。
    """
    stem = title.strip() or source_path.stem
    stem = re.sub(r'[<>:"/\\|?*\r\n\t]+', "_", stem)
    stem = re.sub(r"\s+", "_", stem)
    stem = re.sub(r"_+", "_", stem).strip("._ ")
    if not stem:
        stem = "wiki_" + hashlib.md5(source_path.name.encode()).hexdigest()[:8]
    return f"{stem}_{TODAY.replace('-', '')}.md"


def save_wiki(wiki_content: str, dest_dir: Path, filename: str) -> Path:
    """wiki Markdown をファイルに保存する"""
    dest_dir.mkdir(parents=True, exist_ok=True)
    wiki_path = dest_dir / filename

    # 同名ファイルが存在する場合は _v2, _v3 を付与
    if wiki_path.exists():
        base = wiki_path.stem
        i = 2
        while wiki_path.exists():
            wiki_path = dest_dir / f"{base}_v{i}.md"
            i += 1

    wiki_path.write_text(wiki_content, encoding="utf-8")
    return wiki_path


# ════════════════════════════════════════════════════════
# processed-sources.yaml 更新
# ════════════════════════════════════════════════════════

def compute_md5(file_path: Path) -> str:
    """ファイルの MD5 ハッシュを計算する"""
    h = hashlib.md5()
    with open(file_path, "rb") as f:
        for chunk in iter(lambda: f.read(65536), b""):
            h.update(chunk)
    return h.hexdigest()


def to_personal_relative(abs_path: Path) -> str:
    """絶対パスを 00personal 相対パス（スラッシュ区切り）に変換する"""
    return abs_path.relative_to(PERSONAL_ROOT).as_posix()


def to_kb_relative(abs_path: Path) -> str:
    """絶対パスを KnowledgeBase 相対パス（スラッシュ区切り）に変換する"""
    return abs_path.relative_to(KB_ROOT).as_posix()


def update_processed_sources(source_path: Path, wiki_path: Path) -> None:
    """
    processed-sources.yaml に処理レコードを upsert する。

    以前の実装は常に追記していたが、同じ source_path を再処理すると重複が増える。
    batch-inbox.md の重複防止は source_path をキーにするため、本番保存モードでは
    既存レコードを置き換え、なければ追加する。
    """
    record = {
        "source_path": to_personal_relative(source_path),
        "wiki_path": to_kb_relative(wiki_path),
        "file_hash": compute_md5(source_path),
        "processed_date": TODAY,
        # このスクリプト単体では index-builder を実行しない。
        # 週次メンテナンスまたは後続の index-builder add で true にする前提。
        "index_registered": False,
        "binary_moved": "skipped",
        "binary_destination": to_personal_relative(source_path),
        "source_deleted": False,
    }

    records = []
    if PROCESSED_PATH.exists():
        loaded = yaml.safe_load(PROCESSED_PATH.read_text(encoding="utf-8")) or []
        if isinstance(loaded, list):
            records = loaded

    replaced = False
    for i, existing in enumerate(records):
        if isinstance(existing, dict) and existing.get("source_path") == record["source_path"]:
            records[i] = record
            replaced = True
            break
    if not replaced:
        records.append(record)

    PROCESSED_PATH.parent.mkdir(parents=True, exist_ok=True)
    PROCESSED_PATH.write_text(
        yaml.dump(records, allow_unicode=True, default_flow_style=False, sort_keys=False),
        encoding="utf-8",
    )

    print(f"📋 processed-sources.yaml を更新しました", file=sys.stderr)


def build_frontmatter(analysis: dict, source_path: Path, status: str) -> str:
    """
    Gemini分類結果からFront-matterを作る。

    本文生成はGeminiへ寄せるが、Front-matterは機械的に作る。
    これにより、LLMがYAMLを崩したり必須項目を落とす事故を避ける。
    """
    tags = analysis.get("tags") or []
    fm = {
        "wiki_type": analysis.get("wiki_type", "reference"),
        "title": analysis.get("title") or source_path.stem,
        "aliases": [],
        "created": TODAY,
        "updated": TODAY,
        "source": "internal_doc",
        "source_url": "",
        "tags": tags,
        "scope": infer_scope_from_destination(analysis.get("destination", "")),
        "domain": infer_domain_from_destination(analysis.get("destination", "")),
        "status": status,
        "related": [],
        "source_file": to_personal_relative(source_path),
        "classification_confidence": analysis.get("confidence_score"),
        "classification_method": analysis.get("classification_method", "gemini_schema_hints"),
        "needs_review": bool(analysis.get("needs_review", False)),
    }
    return "---\n" + yaml.dump(fm, allow_unicode=True, sort_keys=False).strip() + "\n---\n\n"


def infer_scope_from_destination(destination: str) -> str:
    """保存先のトップ階層からFront-matterのscopeを推定する。"""
    top = destination.strip("/").split("/", 1)[0] if destination else ""
    return {
        "kyorindo": "kyorindo",
        "tsuruha-hd": "tsuruha-hd",
        "retail": "retail",
        "ai-dx": "industry",
        "vendor": "industry",
        "research": "general",
    }.get(top, "general")


def infer_domain_from_destination(destination: str) -> str:
    """保存先パスからFront-matterのdomainを推定する。"""
    parts = destination.strip("/").split("/") if destination else []
    if not parts:
        return "it-systems"
    if parts[0] == "kyorindo" and len(parts) > 1:
        return parts[1]
    return parts[0]


def save_wiki_from_analysis(source_path: Path, analysis: dict, body: str) -> Path:
    """
    分類結果と本文から、KnowledgeBaseへwikiファイルを保存する。

    needs_review=true の分類は、destination がある場合のみ draft として保存する。
    destination が空の場合は保存先が決められないため、呼び出し側でスキップする。
    """
    destination = analysis.get("destination", "").strip("/")
    if not destination:
        raise ValueError("destination が空のためwiki保存をスキップします")
    status = "draft" if analysis.get("needs_review") else "current"
    dest_dir = KB_ROOT / destination
    filename = suggest_wiki_filename_from_title(analysis.get("title", ""), source_path)
    content = build_frontmatter(analysis, source_path, status) + body.strip() + "\n"
    return save_wiki(content, dest_dir, filename)


# ════════════════════════════════════════════════════════
# メイン処理
# ════════════════════════════════════════════════════════

def main():
    parser = argparse.ArgumentParser(description="Gemini API を使って wiki Markdown を生成する")
    parser.add_argument("file_path", help="処理対象ファイルのパス（xlsx/pptx/pdf）")
    # スタンドアロン用
    parser.add_argument("--dest",    default=None,  help="wiki保存先（KnowledgeBase相対パス）例: kyorindo/it-systems/")
    parser.add_argument("--dry-run", action="store_true", help="ファイル保存・YAML更新をせず出力のみ確認")
    # batch-inbox.md の低コスト経路用。
    # --analyze-only は analyze.md の代替分類、--generate は分類結果を受けた本文生成を担当する。
    parser.add_argument("--analyze-only", action="store_true", help="分類結果JSONのみ出力する（batch-inbox.md低コスト経路用）")
    parser.add_argument("--generate",     action="store_true", help="分類結果に基づいて本文Markdownのみ出力する")
    parser.add_argument("--save-wiki",    action="store_true", help="分類・本文生成・wiki保存・processed-sources更新を行う")
    parser.add_argument("--analysis-json", default="", help="--generate 用の分類結果JSON文字列またはJSONファイルパス")
    parser.add_argument("--analysis-json-stdin", action="store_true", help="--generate 用の分類結果JSONを標準入力から読む")
    parser.add_argument("--schema-path", default=str(SCHEMA_PATH), help="--analyze-only 用 SCHEMA.md パス")
    parser.add_argument("--hints-path",  default=str(HINTS_PATH),  help="--analyze-only 用 classification-hints.md パス")
    parser.add_argument("--emit-usage", action="store_true", help="Gemini usage_metadata と入出力文字数を stderr/JSON に出力する")
    # 互換用（--body-only）
    parser.add_argument("--body-only",      action="store_true", help="本文のみ stdout に出力（write-wiki.md統合用）")
    parser.add_argument("--wiki-type",      default="reference", help="wiki_type（analyze.mdの分類結果）")
    parser.add_argument("--title",          default="",          help="タイトル候補（analyze.mdの分類結果）")
    parser.add_argument("--scope",          default="",          help="スコープ（kyorindo等）")
    parser.add_argument("--domain",         default="",          help="ドメイン（cx/it-systems等）")
    parser.add_argument("--tags",           default="",          help="タグ（カンマ区切り）")
    parser.add_argument("--converted-text", default="",          help="事前抽出済みテキスト（省略時はファイルから抽出）")
    args = parser.parse_args()

    file_path = Path(args.file_path).resolve()

    # ── バリデーション ────────────────────────────────────
    if not file_path.exists():
        print(f"❌ ファイルが見つかりません: {file_path}", file=sys.stderr)
        sys.exit(1)

    ext = file_path.suffix.lower()
    if ext not in (".xlsx", ".pptx", ".pdf", ".md", ".txt"):
        print(f"❌ 未対応の形式です（xlsx/pptx/pdf/md/txt のみ対応）: {ext}", file=sys.stderr)
        sys.exit(1)

    # ── テキスト抽出 ─────────────────────────────────────
    if args.converted_text:
        # 事前抽出済みテキストを使用（convert-binary.md の結果を再利用）
        content = args.converted_text
        if not (args.body_only or args.generate or args.analyze_only or args.save_wiki):
            print(f"📄 事前抽出テキストを使用: {len(content):,} 文字")
    else:
        if not (args.body_only or args.generate or args.analyze_only or args.save_wiki):
            print(f"📄 テキスト抽出中: {file_path.name}")
        try:
            content = extract_text(file_path)
        except Exception as e:
            print(f"❌ テキスト抽出失敗: {e}", file=sys.stderr)
            sys.exit(1)

        if not content.strip():
            content = f"ファイル名: {file_path.name}\nフォルダ: {file_path.parent.name}"
            if not (args.body_only or args.generate or args.analyze_only or args.save_wiki):
                print("⚠️  テキスト抽出不可。ファイル名から内容を推定します。")

        if not (args.body_only or args.generate or args.analyze_only or args.save_wiki):
            print(f"   抽出文字数: {len(content):,} 文字")

    # ════════════════════════════════════════════════════
    # --save-wiki モード
    # 本番用: 分類 → 本文生成 → Front-matter付き保存 → processed-sources更新を実行する。
    # destination が空のものは安全のため保存せず skipped として返す。
    # needs_review=true だが destination があるものは draft として保存する。
    # ════════════════════════════════════════════════════
    if args.save_wiki:
        try:
            analysis, analyze_usage = analyze_document(
                file_path=file_path,
                content=content,
                schema_path=Path(args.schema_path),
                hints_path=Path(args.hints_path),
            )
        except Exception as e:
            print(f"❌ Gemini分類エラー: {e}", file=sys.stderr)
            sys.exit(1)

        result = {
            "status": "skipped",
            "source_path": to_personal_relative(file_path),
            "analysis": analysis,
            "wiki_path": "",
            "usage": {"analyze": analyze_usage, "generate": None},
        }
        if not analysis.get("destination"):
            result["skip_reason"] = "destination_empty"
            sys.stdout.reconfigure(encoding="utf-8")
            print(json.dumps(result, ensure_ascii=False, indent=2))
            return

        try:
            body, generate_usage = generate_wiki_body(
                file_path=file_path,
                content=content,
                wiki_type=analysis.get("wiki_type", args.wiki_type),
                title=analysis.get("title", args.title),
                destination=analysis.get("destination", args.dest or ""),
                scope=args.scope,
                domain=args.domain,
                tags=analysis.get("tags") or [],
            )
            wiki_path = save_wiki_from_analysis(file_path, analysis, body)
            update_processed_sources(file_path, wiki_path)
        except Exception as e:
            print(f"❌ wiki保存エラー: {e}", file=sys.stderr)
            sys.exit(1)

        result["status"] = "success"
        result["wiki_path"] = to_kb_relative(wiki_path)
        result["usage"]["generate"] = generate_usage
        sys.stdout.reconfigure(encoding="utf-8")
        print(json.dumps(result, ensure_ascii=False, indent=2))
        return

    # ════════════════════════════════════════════════════
    # --analyze-only モード
    # analyze.md の高コスト分類を Python/Gemini 側へ逃がすための新しい正本候補。
    # stdout は機械処理しやすいJSONだけにし、usageログは --emit-usage 時にJSON内へ含める。
    # ════════════════════════════════════════════════════
    if args.analyze_only:
        try:
            analysis, usage = analyze_document(
                file_path=file_path,
                content=content,
                schema_path=Path(args.schema_path),
                hints_path=Path(args.hints_path),
            )
        except Exception as e:
            print(f"❌ Gemini分類エラー: {e}", file=sys.stderr)
            sys.exit(1)
        if args.emit_usage:
            analysis["usage"] = usage
        sys.stdout.reconfigure(encoding="utf-8")
        print(json.dumps(analysis, ensure_ascii=False, indent=2))
        return

    # ════════════════════════════════════════════════════
    # --generate モード
    # --analyze-only の分類結果を受け取り、本文生成だけを行う。
    # write-wiki.md互換の --body-only と同じ本文生成関数を使うが、
    # 将来 batch-inbox.md から直接呼ぶために明示モードとして分けている。
    # ════════════════════════════════════════════════════
    if args.generate:
        analysis = _load_analysis_from_args(args)
        tags = analysis.get("tags")
        if tags is None:
            tags = [t.strip() for t in args.tags.split(",") if t.strip()]
        try:
            body, usage = generate_wiki_body(
                file_path=file_path,
                content=content,
                wiki_type=analysis.get("wiki_type", args.wiki_type),
                title=analysis.get("title", args.title),
                destination=analysis.get("destination", args.dest or ""),
                scope=args.scope,
                domain=args.domain,
                tags=tags,
            )
        except Exception as e:
            print(f"❌ Gemini本文生成エラー: {e}", file=sys.stderr)
            sys.exit(1)
        if args.emit_usage:
            _print_usage(usage)
        sys.stdout.reconfigure(encoding="utf-8")
        print(body)
        return

    # ════════════════════════════════════════════════════
    # --body-only モード（write-wiki.md スキル統合用）
    # Front-matterなし・本文のみを stdout に出力して終了
    # ════════════════════════════════════════════════════
    if args.body_only:
        tags = [t.strip() for t in args.tags.split(",") if t.strip()]
        try:
            body, usage = generate_wiki_body(
                file_path=file_path,
                content=content,
                wiki_type=args.wiki_type,
                title=args.title,
                destination=args.dest or "",
                scope=args.scope,
                domain=args.domain,
                tags=tags,
            )
        except Exception as e:
            print(f"❌ Gemini API エラー: {e}", file=sys.stderr)
            sys.exit(1)
        if args.emit_usage:
            _print_usage(usage)
        # 本文のみを stdout に出力（write-wiki.md が受け取る）
        sys.stdout.reconfigure(encoding="utf-8")
        print(body)
        return

    # ════════════════════════════════════════════════════
    # スタンドアロンモード（通常 / dry-run）
    # ════════════════════════════════════════════════════
    try:
        wiki_content, usage = generate_wiki(file_path, content)
    except Exception as e:
        print(f"❌ Gemini API エラー: {e}")
        sys.exit(1)
    if args.emit_usage:
        _print_usage(usage)

    # ── dry-run モード ───────────────────────────────────
    if args.dry_run:
        print("\n" + "═" * 60)
        print("【DRY-RUN】生成されたwiki内容（保存はしません）:")
        print("═" * 60)
        print(wiki_content)
        return

    # ── 保存先の決定 ─────────────────────────────────────
    if args.dest:
        dest_dir = KB_ROOT / args.dest.strip("/")
    else:
        # デフォルト: KnowledgeBase/_inbox_gemini/ に保存（手動で移動）
        dest_dir = KB_ROOT / "_inbox_gemini"
        print(f"⚠️  --dest 未指定のため {dest_dir} に保存します")

    filename = suggest_wiki_filename(file_path)

    # ── ファイル保存 ─────────────────────────────────────
    wiki_path = save_wiki(wiki_content, dest_dir, filename)
    print(f"✅ wiki保存完了: {wiki_path}")

    # ── processed-sources.yaml 更新 ─────────────────────
    try:
        update_processed_sources(file_path, wiki_path)
    except Exception as e:
        print(f"⚠️  YAML更新失敗（手動で追記してください）: {e}")

    print("\n🎉 完了!")
    print(f"   ソース : {to_personal_relative(file_path)}")
    print(f"   wiki   : {to_kb_relative(wiki_path)}")


if __name__ == "__main__":
    main()
